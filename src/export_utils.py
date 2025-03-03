import os
import base64
from pathlib import Path
import markdown2
import docx
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import pdfkit
from weasyprint import HTML
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from pptx.dml.color import RGBColor
import tempfile
import re

class ExportUtils:
    """Utility class for exporting research reports to various formats"""
    
    @staticmethod
    def markdown_to_html(markdown_content, css=None):
        """Convert markdown to HTML with optional CSS styling"""
        # Basic CSS for styling if none provided
        if css is None:
            css = """
            <style>
                body {
                    font-family: Arial, sans-serif;
                    line-height: 1.6;
                    margin: 40px;
                    max-width: 900px;
                    margin: 0 auto;
                    padding: 20px;
                    color: #333;
                }
                h1 { font-size: 28px; margin-top: 20px; color: #1a73e8; }
                h2 { font-size: 24px; margin-top: 18px; color: #1a73e8; border-bottom: 1px solid #eee; padding-bottom: 5px; }
                h3 { font-size: 20px; margin-top: 16px; color: #1a73e8; }
                p { margin-bottom: 16px; }
                a { color: #1a73e8; text-decoration: none; }
                a:hover { text-decoration: underline; }
                code { background-color: #f5f5f5; padding: 2px 5px; border-radius: 3px; font-family: monospace; }
                pre { background-color: #f5f5f5; padding: 15px; border-radius: 5px; overflow-x: auto; }
                blockquote { border-left: 4px solid #ddd; padding-left: 15px; color: #666; }
                img { max-width: 100%; height: auto; }
                table { border-collapse: collapse; width: 100%; margin-bottom: 20px; }
                th, td { border: 1px solid #ddd; padding: 8px 12px; text-align: left; }
                th { background-color: #f5f5f5; }
                tr:nth-child(even) { background-color: #f9f9f9; }
            </style>
            """
        
        # Convert markdown to HTML
        html_content = markdown2.markdown(
            markdown_content,
            extras=["tables", "code-friendly", "fenced-code-blocks", "header-ids", "footnotes"]
        )
        
        # Combine with a full HTML document structure
        full_html = f"""<!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Research Report</title>
            {css}
        </head>
        <body>
            {html_content}
        </body>
        </html>
        """
        
        return full_html

    @staticmethod
    def export_to_html(markdown_content, output_path):
        """Export markdown content to a standalone HTML file"""
        html_content = ExportUtils.markdown_to_html(markdown_content)
        
        # Write HTML to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path

    @staticmethod
    def export_to_pdf(markdown_content, output_path, method='weasyprint'):
        """
        Export markdown content to PDF
        
        Args:
            markdown_content: The markdown content to convert
            output_path: Path to save the PDF
            method: 'weasyprint' or 'pdfkit' (default: weasyprint)
        """
        html_content = ExportUtils.markdown_to_html(markdown_content)
        
        if method == 'pdfkit':
            try:
                # Using pdfkit (requires wkhtmltopdf installed)
                options = {
                    'page-size': 'A4',
                    'margin-top': '20mm',
                    'margin-right': '20mm',
                    'margin-bottom': '20mm',
                    'margin-left': '20mm',
                    'encoding': 'UTF-8',
                }
                pdfkit.from_string(html_content, output_path, options=options)
            except Exception as e:
                # Fallback to weasyprint if pdfkit fails
                print(f"Error with pdfkit: {str(e)}. Falling back to weasyprint.")
                method = 'weasyprint'
        
        if method == 'weasyprint':
            # Using WeasyPrint
            with tempfile.NamedTemporaryFile(suffix='.html', mode='w', encoding='utf-8', delete=False) as f:
                f.write(html_content)
                temp_html_path = f.name
            
            # Convert to PDF
            HTML(filename=temp_html_path).write_pdf(output_path)
            
            # Clean up temporary file
            try:
                os.remove(temp_html_path)
            except:
                pass
        
        return output_path

    @staticmethod
    def export_to_docx(markdown_content, output_path):
        """Export markdown content to DOCX format"""
        doc = docx.Document()
        
        # Set document properties
        doc.core_properties.title = "Research Report"
        
        # Parse markdown sections
        lines = markdown_content.split('\n')
        
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # Handle headers
            if line.startswith('# '):
                # Title (H1)
                title = line.replace('# ', '')
                heading = doc.add_heading(title, level=0)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif line.startswith('## '):
                # Section title (H2)
                heading = doc.add_heading(line.replace('## ', ''), level=1)
            elif line.startswith('### '):
                # Subsection (H3)
                heading = doc.add_heading(line.replace('### ', ''), level=2)
            elif line.startswith('#### '):
                heading = doc.add_heading(line.replace('#### ', ''), level=3)
            elif line.startswith('- ') or line.startswith('* '):
                # Bullet point
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line.startswith('1. ') or (line.startswith(str(i)) and '. ' in line):
                # Numbered list
                doc.add_paragraph(line[line.find('.')+2:], style='List Number')
            elif line.startswith('```'):
                # Code block
                code_block = []
                i += 1
                while i < len(lines) and not lines[i].startswith('```'):
                    code_block.append(lines[i])
                    i += 1
                if code_block:
                    code_para = doc.add_paragraph()
                    code_text = code_para.add_run('\n'.join(code_block))
                    code_text.font.name = 'Courier New'
                    code_text.font.size = Pt(9)
            elif line == '':
                # Empty line
                if i > 0 and i < len(lines) - 1 and lines[i-1].strip() != '' and lines[i+1].strip() != '':
                    doc.add_paragraph()
            else:
                # Regular paragraph
                if line:
                    # Process bold and italic markdown
                    line = re.sub(r'\*\*(.+?)\*\*', lambda m: m.group(1), line)  # Bold
                    line = re.sub(r'_(.+?)_', lambda m: m.group(1), line)        # Italic
                    line = re.sub(r'\[(.+?)\]\((.+?)\)', lambda m: f"{m.group(1)}", line)  # Links
                    
                    doc.add_paragraph(line)
            
            i += 1
        
        # Save the document
        doc.save(output_path)
        return output_path

    @staticmethod
    def export_to_presentation(markdown_content, output_path):
        """Export markdown content to PowerPoint presentation format"""
        prs = Presentation()
        
        # Set presentation properties
        title_slide_layout = prs.slide_layouts[0]  # Title slide
        content_slide_layout = prs.slide_layouts[1]  # Title and content
        
        # Parse markdown for sections
        sections = markdown_content.split('## ')
        
        # Create title slide from the H1 header
        title_match = re.search(r'# (.*?)$', sections[0], re.MULTILINE)
        if title_match:
            title_text = title_match.group(1).strip()
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = title_text
            
            # Add subtitle if there's content before the first H2
            subtitle_text = sections[0].replace(f"# {title_text}", "").strip()
            if subtitle_text:
                subtitle = slide.placeholders[1]
                subtitle.text = "Deep Research Report"
        
        # Process each section (starting from 1 to skip the title part)
        for section in sections[1:]:
            lines = section.split('\n')
            if not lines:
                continue
                
            # Get section title (first line is the H2 title)
            section_title = lines[0].strip()
            
            # Create a new slide for this section
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            title.text = section_title
            
            # Add content
            content = "\n".join(lines[1:])
            
            # Process content - simplify for presentation
            # Remove code blocks, simplify lists
            content = re.sub(r'```.*?```', '', content, flags=re.DOTALL)
            content = re.sub(r'!\[.*?\]\(.*?\)', '', content)  # Remove images
            
            # Extract key points
            content = content.strip()
            if content:
                body_shape = slide.shapes.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.text = ""
                
                # Split into paragraphs
                paragraphs = content.split('\n\n')
                for i, para in enumerate(paragraphs):
                    if i >= 6:  # Limit to 6 points per slide
                        # If we have more content, create a new slide
                        slide = prs.slides.add_slide(content_slide_layout)
                        title = slide.shapes.title
                        title.text = section_title + " (continued)"
                        body_shape = slide.shapes.placeholders[1]
                        text_frame = body_shape.text_frame
                        text_frame.text = ""
                        i = 0
                    
                    if para.strip():
                        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                        
                        # Convert bullet points and simplify
                        para = para.strip()
                        if para.startswith('- ') or para.startswith('* '):
                            p.text = para[2:].strip()
                            p.level = 0
                        elif para.startswith('  - ') or para.startswith('  * '):
                            p.text = para[4:].strip()
                            p.level = 1
                        else:
                            p.text = para
        
        # Save the presentation
        prs.save(output_path)
        return output_path

    @staticmethod
    def get_file_download_link(file_path, label="Download File"):
        """Generate a download link for a file"""
        with open(file_path, 'rb') as f:
            data = f.read()
        
        b64 = base64.b64encode(data).decode()
        file_name = os.path.basename(file_path)
        mime_type = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.html': 'text/html',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }.get(os.path.splitext(file_path)[1], 'application/octet-stream')
        
        href = f'data:{mime_type};base64,{b64}'
        download_link = f'''
        <a href="{href}" download="{file_name}" 
           style="text-decoration: none; color: white;">
            <button style="background-color: #4CAF50; color: white; padding: 8px 16px;
                    border: none; border-radius: 4px; cursor: pointer;">
                {label}
            </button>
        </a>
        '''
        return download_link